import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color palette
    C_NAVY_DARK  = RGBColor(12, 38, 48)     # #0c2630
    C_NAVY_LIGHT = RGBColor(18, 60, 74)     # #123c4a
    C_CYAN       = RGBColor(0, 220, 245)    # #00dcf5 (Bright Cyan from SeniorClub)
    C_CYAN_DARK  = RGBColor(0, 160, 195)    # #00a0c3
    C_BG_LIGHT   = RGBColor(248, 250, 252)  # #f8fafc
    C_WHITE      = RGBColor(255, 255, 255)  # #ffffff
    C_CARD_BG    = RGBColor(255, 255, 255)  # #ffffff
    C_TEXT_DARK  = RGBColor(15, 23, 42)     # #0f172a
    C_TEXT_MUTED = RGBColor(100, 116, 139)  # #64748b
    C_BORDER     = RGBColor(226, 232, 240)  # #e2e8f0

    blank_layout = prs.slide_layouts[6]

    def set_bg(slide, color=C_BG_LIGHT):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_header_light(slide, badge_text, title_text, subtitle_text=""):
        # Header top bar accent
        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.1))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = C_CYAN
        top_bar.line.fill.background()

        # Badge Pill
        badge_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(0.4), Inches(2.8), Inches(0.35))
        badge_bg.fill.solid()
        badge_bg.fill.fore_color.rgb = C_NAVY_LIGHT
        badge_bg.line.fill.background()
        tf_b = badge_bg.text_frame
        tf_b.word_wrap = True
        p_b = tf_b.paragraphs[0]
        p_b.text = badge_text.upper()
        p_b.font.size = Pt(10)
        p_b.font.bold = True
        p_b.font.color.rgb = C_CYAN
        p_b.alignment = PP_ALIGN.CENTER

        # Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.7))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(26)
        p2.font.bold = True
        p2.font.color.rgb = C_TEXT_DARK

        # Subtitle
        if subtitle_text:
            txBox3 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5))
            tf3 = txBox3.text_frame
            tf3.word_wrap = True
            p3 = tf3.paragraphs[0]
            p3.text = subtitle_text
            p3.font.size = Pt(13)
            p3.font.color.rgb = C_TEXT_MUTED

    # ---------------------------------------------------------
    # SLIDE 1: Cover (Dark Navy Header + Cyan Accents)
    # ---------------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1, C_NAVY_DARK)

    # Cyan Accent Bar
    bar1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
    bar1.fill.solid()
    bar1.fill.fore_color.rgb = C_CYAN
    bar1.line.fill.background()

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "●  BENEFICIO CORPORATIVO DE BIENESTAR"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_CYAN

    p2 = tf.add_paragraph()
    p2.text = "CompanyCare"
    p2.font.size = Pt(56)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(6)

    p3 = tf.add_paragraph()
    p3.text = "Plataforma de cuidado de adultos mayores para equipos que importan.\nConectamos colaboradores con expertos en gerontología y convenios exclusivos."
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_before = Pt(16)

    # ---------------------------------------------------------
    # SLIDE 2: El Problema
    # ---------------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2)
    add_header_light(slide2, "El Problema", "¿Por qué CompanyCare?", "Millones de colaboradores cuidan a un familiar adulto mayor mientras trabajan.")

    cards_data = [
        ("Estrés del Cuidador", "73%", "de los cuidadores reportan altos niveles de estrés que afectan su salud mental y desempeño laboral."),
        ("Ausentismo Laboral", "6.6 hrs", "semanales pierden los colaboradores en promedio gestionando el cuidado de su familiar mayor."),
        ("Falta de Orientación", "80%", "no sabe dónde buscar residencias, cuidadores o apoyo profesional de confianza.")
    ]

    for i, (ctitle, stat, cdesc) in enumerate(cards_data):
        left = Inches(0.8 + i * 3.9)
        shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.3), Inches(3.6), Inches(4.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_BORDER

        tf = shape.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = ctitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_LIGHT

        p_stat = tf.add_paragraph()
        p_stat.text = stat
        p_stat.font.size = Pt(36)
        p_stat.font.bold = True
        p_stat.font.color.rgb = C_CYAN_DARK
        p_stat.space_before = Pt(10)

        p2 = tf.add_paragraph()
        p2.text = cdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(8)

    # ---------------------------------------------------------
    # SLIDE 3: Los Roles (Vibrant High-Contrast Layout)
    # ---------------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3)
    add_header_light(slide3, "Usuarios", "¿Quiénes usan la plataforma?", "CompanyCare conecta 4 roles en un ecosistema integrado.")

    roles = [
        ("Colaborador", "Empleado que cuida a un adulto mayor. Accede al beneficio corporativo.", "Dashboard, Asesoría, Proveedores, Recursos, SeniorClub", C_NAVY_LIGHT),
        ("Care Expert", "Especialista en gerontología y trabajo social que asesora a los colaboradores.", "Inbox de casos, Chat clínico, Evolución del familiar", C_CYAN_DARK),
        ("RRHH / Admin Empresa", "Administrador de la empresa que gestiona empleados y monitorea casos.", "Panel de empresa, Monitoreo de bienestar, Vouchers", C_NAVY_LIGHT),
        ("Super Admin", "Administrador global de la plataforma CompanyCare / Senior Advisor.", "ERP interno, Gestión multi-tenant, Facturación", C_CYAN_DARK)
    ]

    for i, (rname, rdesc, racces, tag_col) in enumerate(roles):
        top = Inches(2.2 + i * 1.18)
        shape = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.05))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_BORDER

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = rname
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = tag_col

        p2 = tf.add_paragraph()
        p2.text = f"{rdesc}   |   Acceso: {racces}"
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(3)

    # ---------------------------------------------------------
    # SLIDE 4: Flujo General
    # ---------------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4)
    add_header_light(slide4, "Flujo General", "¿Cómo funciona paso a paso?", "Desde la activación del beneficio hasta el monitoreo continuo de RRHH.")

    steps = [
        ("PASO 1", "Empresa Activa", "La empresa se registra, elige su plan y el Admin invita a sus colaboradores al portal."),
        ("PASO 2", "Colaborador Ingresa", "El empleado accede y completa su ficha de cuidado con los datos de su familiar mayor."),
        ("PASO 3", "Care Expert Asesora", "El experto toma el caso, realiza videollamada, chat o llamada y registra la evolución."),
        ("PASO 4", "RRHH Monitorea", "El admin de empresa ve en tiempo real el estado de bienestar de todos sus colaboradores.")
    ]

    for i, (tag, stitle, sdesc) in enumerate(steps):
        left = Inches(0.8 + i * 2.95)
        shape = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.3), Inches(2.75), Inches(4.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_BORDER

        tf = shape.text_frame
        tf.word_wrap = True
        
        p_tag = tf.paragraphs[0]
        p_tag.text = tag
        p_tag.font.size = Pt(11)
        p_tag.font.bold = True
        p_tag.font.color.rgb = C_CYAN_DARK

        p = tf.add_paragraph()
        p.text = stitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_LIGHT
        p.space_before = Pt(6)

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(10)

    # ---------------------------------------------------------
    # SLIDE 5: Dashboard Colaborador
    # ---------------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5)
    add_header_light(slide5, "Paso 2 · Colaborador", "Centro de Control del Colaborador", "El empleado ve de un vistazo el estado de su familiar y sus solicitudes activas.")

    box5 = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    box5.fill.solid()
    box5.fill.fore_color.rgb = C_CARD_BG
    box5.line.color.rgb = C_BORDER
    tf5 = box5.text_frame
    tf5.word_wrap = True

    p = tf5.paragraphs[0]
    p.text = "Funcionalidades Clave para el Colaborador:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_LIGHT

    items5 = [
        "✔ Estado de salud del familiar actualizado en tiempo real por el Care Expert ('Estable', 'Mejorando', 'Sin cambios').",
        "✔ Mi Ficha de Cuidado: Formulario único para registrar condiciones médicas, dependencia, previsión y red de apoyo.",
        "✔ Historial de evolución completo con notas del especialista y fechas de próximo seguimiento.",
        "✔ Solicitudes activas y agendamiento de consultas por Videollamada (Google Meet), Llamada o Chat.",
        "✔ Acceso a SeniorClub: Descuentos exclusivos en salud, telemedicina, ayudas técnicas y vestuario adaptado."
    ]
    for itm in items5:
        p_item = tf5.add_paragraph()
        p_item.text = itm
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = C_TEXT_MUTED
        p_item.space_before = Pt(12)

    # ---------------------------------------------------------
    # SLIDE 6: Asesoría Care Expert
    # ---------------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6)
    add_header_light(slide6, "Paso 3 · Asesoría", "Conexión con Care Experts", "El colaborador agenda una consulta con un especialista por el canal que prefiera.")

    channels = [
        ("Videollamada", "Google Meet automático", "Reunión cara a cara por videoconferencia generada instantáneamente en la plataforma."),
        ("Llamada Telefónica", "Directo por voz", "Contacto telefónico inmediato para resolver dudas urgentes sin necesidad de conexión a internet."),
        ("Chat en Tiempo Real", "Orientación continua", "Atención escrita fluida con historial de conversación y envío de documentos relevantes.")
    ]

    for i, (ch_title, ch_tag, ch_desc) in enumerate(channels):
        left = Inches(0.8 + i * 3.9)
        shape = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.3), Inches(3.6), Inches(4.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_BORDER

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ch_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_LIGHT

        p2 = tf.add_paragraph()
        p2.text = f"● {ch_tag}"
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = C_CYAN_DARK
        p2.space_before = Pt(6)

        p3 = tf.add_paragraph()
        p3.text = ch_desc
        p3.font.size = Pt(12)
        p3.font.color.rgb = C_TEXT_MUTED
        p3.space_before = Pt(10)

    # ---------------------------------------------------------
    # SLIDE 7: Panel Care Expert
    # ---------------------------------------------------------
    slide7 = prs.slides.add_slide(blank_layout)
    set_bg(slide7)
    add_header_light(slide7, "Care Expert", "Espacio Clínico del Care Expert", "El experto gestiona todos sus casos desde una suite profesional unificada.")

    box7 = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    box7.fill.solid()
    box7.fill.fore_color.rgb = C_CARD_BG
    box7.line.color.rgb = C_BORDER
    tf7 = box7.text_frame
    tf7.word_wrap = True

    p = tf7.paragraphs[0]
    p.text = "Herramientas de Trabajo del Care Expert:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_LIGHT

    items7 = [
        "✔ Inbox de Casos Activos con indicadores de SLA (>24h crítico, >6h advertencia).",
        "✔ Chat directo con el cliente + pestañas de Notas Internas Confidenciales.",
        "✔ Asignación instantánea de casos sin asignar ('Tomar caso') con auto-asociación de citas.",
        "✔ Formulario de Evolución del Familiar: Estado (Estable, Mejorando, Empeorando), Prioridad y fecha de próximo seguimiento.",
        "✔ Ficha contextual completa del colaborador y su adulto mayor (RUT, diagnóstico, nivel de dependencia, escaleras/rampas)."
    ]
    for itm in items7:
        p_item = tf7.add_paragraph()
        p_item.text = itm
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = C_TEXT_MUTED
        p_item.space_before = Pt(12)

    # ---------------------------------------------------------
    # SLIDE 8: Panel Empresa
    # ---------------------------------------------------------
    slide8 = prs.slides.add_slide(blank_layout)
    set_bg(slide8)
    add_header_light(slide8, "Empresa", "Panel de Administración Empresa", "Gestión integral de empleados, vouchers, suscripciones y fichas de cuidado.")

    modules = [
        ("Empleados", "Invitar por email, asignar roles (empleado, manager, RRHH) y desvincular."),
        ("Vouchers", "Crear cupones de descuento en porcentaje o monto fijo para servicios de cuidado."),
        ("Suscripciones", "Planes Plataforma y Acompañamiento integrados con Mercado Pago."),
        ("Fichas de Cuidado", "Revisar las fichas completadas por los colaboradores de la empresa.")
    ]

    for i, (mtitle, mdesc) in enumerate(modules):
        left = Inches(0.8 + i * 2.95)
        shape = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.3), Inches(2.75), Inches(4.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_BORDER

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = mtitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_LIGHT

        p2 = tf.add_paragraph()
        p2.text = mdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(12)

    # ---------------------------------------------------------
    # SLIDE 9: Monitoreo RRHH
    # ---------------------------------------------------------
    slide9 = prs.slides.add_slide(blank_layout)
    set_bg(slide9)
    add_header_light(slide9, "Paso 4 · RRHH", "Monitoreo de Bienestar en Tiempo Real", "RRHH sigue el estado de salud de los familiares en un dashboard aislado por empresa.")

    box9 = slide9.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(11.7), Inches(4.5))
    box9.fill.solid()
    box9.fill.fore_color.rgb = C_CARD_BG
    box9.line.color.rgb = C_CYAN_DARK
    tf9 = box9.text_frame
    tf9.word_wrap = True

    p = tf9.paragraphs[0]
    p.text = "Características del Panel de Monitoreo RRHH:"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_NAVY_LIGHT

    items9 = [
        "🔒 Aislamiento estricto por empresa: RRHH solo ve colaboradores de su propia empresa (RLS en Supabase).",
        "🚨 Banner de Alerta Urgente para visualizar casos marcados con prioridad urgente.",
        "📊 Filtros rápidos por estado de salud del familiar ('Estable', 'Mejorando', 'Sin cambios', 'Atención req.', 'Empeorando').",
        "🔎 Filas expandibles con detalle completo: última nota del Care Expert, canal, experto asignado y fecha de próximo seguimiento.",
        "⚡ Búsqueda instantánea por nombre de colaborador, tema o contenido de notas."
    ]
    for itm in items9:
        p_item = tf9.add_paragraph()
        p_item.text = itm
        p_item.font.size = Pt(13)
        p_item.font.color.rgb = C_TEXT_MUTED
        p_item.space_before = Pt(12)

    # ---------------------------------------------------------
    # SLIDE 10: Seguridad
    # ---------------------------------------------------------
    slide10 = prs.slides.add_slide(blank_layout)
    set_bg(slide10)
    add_header_light(slide10, "Seguridad", "Seguridad y Privacidad de Datos", "Múltiples capas que garantizan que cada empresa solo acceda a sus propios datos.")

    sec_layers = [
        ("Row Level Security (RLS)", "Políticas en Supabase que aíslan filas a nivel de base de datos usando company_id. Imposible leer datos de otra empresa."),
        ("Filtro por Membresía", "La app valida company_members antes de consultar care_requests, trayendo solo colaboradores de la misma empresa."),
        ("Guards de Autenticación", "Control de acceso por roles (authGuard, internalAdminGuard) que restringen rutas según permisos.")
    ]

    for i, (stitle, sdesc) in enumerate(sec_layers):
        left = Inches(0.8 + i * 3.9)
        shape = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.3), Inches(3.6), Inches(4.4))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_BORDER

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_LIGHT

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(12)

    # ---------------------------------------------------------
    # SLIDE 11: SeniorClub Intro (Cyan Banner Highlight)
    # ---------------------------------------------------------
    slide11 = prs.slides.add_slide(blank_layout)
    set_bg(slide11)
    add_header_light(slide11, "Programa de Beneficios", "SeniorClub — Convenios Exclusivos", "Convenios exclusivos para el bienestar de nuestros mayores con aliados verificados.")

    # Cyan Banner Feature Box
    cyan_banner = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.1))
    cyan_banner.fill.solid()
    cyan_banner.fill.fore_color.rgb = C_NAVY_DARK
    cyan_banner.line.color.rgb = C_CYAN
    tf_cb = cyan_banner.text_frame
    tf_cb.word_wrap = True
    p_cb = tf_cb.paragraphs[0]
    p_cb.text = "SeniorClub  |  Convenios exclusivos para el bienestar de nuestros mayores"
    p_cb.font.size = Pt(16)
    p_cb.font.bold = True
    p_cb.font.color.rgb = C_CYAN
    p_cb_sub = tf_cb.add_paragraph()
    p_cb_sub.text = "Accede a descuentos y servicios preferenciales con nuestros aliados verificados en Chile."
    p_cb_sub.font.size = Pt(12)
    p_cb_sub.font.color.rgb = C_WHITE
    p_cb_sub.space_before = Pt(4)

    club_stats = [
        ("8+ Aliados Verificados", "Empresas seleccionadas que ofrecen productos y servicios para adultos mayores y cuidadores."),
        ("10% a 30% Descuento", "Descuentos preferenciales en salud, psicogerontología, telemedicina, vestuario adaptado y hogar."),
        ("Obtención en 1 Clic", "Los colaboradores obtienen su código de descuento directamente desde la plataforma con un solo clic.")
    ]

    for i, (stitle, sdesc) in enumerate(club_stats):
        left = Inches(0.8 + i * 3.9)
        shape = slide11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(3.5), Inches(3.6), Inches(3.3))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_BORDER

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = stitle
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_LIGHT

        p2 = tf.add_paragraph()
        p2.text = sdesc
        p2.font.size = Pt(12)
        p2.font.color.rgb = C_TEXT_MUTED
        p2.space_before = Pt(10)

    # ---------------------------------------------------------
    # SLIDE 12: SeniorClub Convenios Detalle (Vibrant Grid)
    # ---------------------------------------------------------
    slide12 = prs.slides.add_slide(blank_layout)
    set_bg(slide12)
    add_header_light(slide12, "Convenios SeniorClub", "Nuestros Convenios Activos", "Descuentos y beneficios preferenciales para la comunidad de SeniorAdvisor y CompanyCare.")

    convenios = [
        ("ListaMente", "10% Dcto. Exclusivo", "Estimulación cognitiva y salud mental. Cupón en checkout de listamente.com."),
        ("Psicóloga B. Alvarado", "15% Dcto. ($30.000)", "Psicogerontología y apoyo en duelo o sobrecarga del cuidador. Ref $40.000."),
        ("Quimun", "30% Dcto. / 3 Meses", "Software para residencias. Carga masiva, costo $0 implementación y soporte."),
        ("Promsa", "15% Dcto. Toda la Web", "Ayudas técnicas, movilidad y confort integral para pacientes y cuidadores."),
        ("Quida (Planes Full)", "10% Dcto. (Viña del Mar)", "Monitoreo discreto con sensores en el hogar. Sin cámaras, con dignidad."),
        ("Bilbi Vístete Fácil", "20% Dcto. Todo Catálogo", "Vestuario adaptado diseñado para facilitar el vestir diario en personas mayores."),
        ("Ducha Segura", "15% Dcto. Adaptación", "Modificación rápida de tinas a ducha a nivel de piso para prevenir caídas."),
        ("Help Rescate Médicos", "Desde $8.336/mes", "Rescate médico 24/7, telemedicina y orientación. RM, Valparaíso, Biobío.")
    ]

    for i, (cname, cdesc_badge, cdesc_txt) in enumerate(convenios):
        row = i // 4
        col = i % 4
        left = Inches(0.8 + col * 2.95)
        top = Inches(2.2 + row * 2.3)

        shape = slide12.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.75), Inches(2.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_CARD_BG
        shape.line.color.rgb = C_BORDER

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cname
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_NAVY_LIGHT

        p2 = tf.add_paragraph()
        p2.text = cdesc_badge
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = C_CYAN_DARK
        p2.space_before = Pt(2)

        p3 = tf.add_paragraph()
        p3.text = cdesc_txt
        p3.font.size = Pt(9.5)
        p3.font.color.rgb = C_TEXT_MUTED
        p3.space_before = Pt(4)

    # ---------------------------------------------------------
    # SLIDE 13: Cierre (Dark Navy Footer Banner)
    # ---------------------------------------------------------
    slide13 = prs.slides.add_slide(blank_layout)
    set_bg(slide13, C_NAVY_DARK)

    bar13 = slide13.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.2), Inches(7.5))
    bar13.fill.solid()
    bar13.fill.fore_color.rgb = C_CYAN
    bar13.line.fill.background()

    txBox = slide13.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.3), Inches(4.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "●  RESUMEN Y PRÓXIMOS PASOS"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = C_CYAN

    p2 = tf.add_paragraph()
    p2.text = "CompanyCare"
    p2.font.size = Pt(48)
    p2.font.bold = True
    p2.font.color.rgb = C_WHITE
    p2.space_before = Pt(6)

    p3 = tf.add_paragraph()
    p3.text = "Una plataforma completa que conecta empresas, colaboradores y expertos en gerontología\npara transformar el cuidado de adultos mayores en un beneficio corporativo medible y seguro."
    p3.font.size = Pt(15)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_before = Pt(14)

    p4 = tf.add_paragraph()
    p4.text = "Contacto: contacto@companycare.cl  |  www.companycare.cl"
    p4.font.size = Pt(14)
    p4.font.bold = True
    p4.font.color.rgb = C_CYAN
    p4.space_before = Pt(24)

    out_path = os.path.join(os.path.dirname(__file__), "presentacion-companycare.pptx")
    prs.save(out_path)
    print(f"Presentation updated successfully with SeniorClub Cyan theme to: {out_path}")

if __name__ == "__main__":
    create_deck()
